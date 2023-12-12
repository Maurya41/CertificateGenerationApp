import os;
import requests;

import jpgconverter;
import sendingmail;

from pptx import Presentation
from pptx.util import Inches, Pt


def replaceText(slide,old_text,new_text):
    
    name_to_replace = old_text;
    
    for shape in slide.shapes:
        if shape.has_text_frame:
            if shape.text_frame.text == name_to_replace:
                text_frame = shape.text_frame;
                
                para = text_frame.paragraphs[0];
                font = para.runs[0].font;
                
                # Get the existing formatting
                font_size = font.size;
                font_bold = font.bold;  
                font_italic = font.italic
                font_color = font.color.rgb

                # Clear the existing text
                para.clear();

                new_run = para.add_run()
                new_run.text = new_text;

                # Apply existing formatting to the new run
                new_run.font.size = font_size
                new_run.font.bold = font_bold
                new_run.font.italic = font_italic
                new_run.font.color.rgb = font_color


def changeName(slide,new_name):
    name_to_replace = "Name";
    replaceText(slide,name_to_replace,new_name);
    

def changeBody(slide,new_body):
    name_to_replace = "Lorem ipsum dolor sit amet, consectetur adipiscing elit. Sed imperdiet fermentum fermentum. Aliquam vitae ante arcu. Pellentesque tempor magna eu ipsum hendrerit fermentum.";
    replaceText(slide,name_to_replace,new_body);


def changeAuth(slide,new_name_1,new_name_2):
    name_to_replace = "SAMIRA";
    replaceText(slide,name_to_replace,new_name_1);
    name_to_replace = "OLIVIA";
    replaceText(slide,name_to_replace,new_name_2);


def getPresentationAndSlide():
    presentation_path = "./assets/Home.pptx";
    presentation = Presentation(presentation_path);
    slide = presentation.slides[0];
    return presentation,slide;

def extendedZip(buffer):
    tmp = [_ for _ in buffer if isinstance(_,list)];
    if len(tmp) > 0:
        length = len(tmp[0]);
        for idx,_ in enumerate(tmp):
            if not len(_) == length:
                raise Exception("Fields in the files not have \
                                same length");

        zipped = [];
        for _ in buffer:
            if not isinstance(_,list):
                zipped.append([_] * length);
            else:
                zipped.append(_);

        return zip(zipped[0],zipped[1],zipped[2],zipped[3]);
    else:
        return zip([buffer[0]],[buffer[1]],[buffer[2]],[buffer[3]]);


def upload_certificate(file_path):
    url = "http://34.142.122.214:5000/upload_certificate";
    
    with open(file_path,'rb') as file:
        files = {
            'file' : (file_path,file)
        };

        response = requests.post(
            url,
            files=files
        )

        if response.status_code == 200:
            download_url = response.content;
            print(download_url);
            return download_url;
        else :
            return None;

def generateModifiedCeritificate(buffer):
    tmp = [_[1] for _ in buffer];

    for name,body,auth,mail in extendedZip(tmp):
        presentation,slide = getPresentationAndSlide();
        changeName(slide,name);
        changeBody(slide,body);
        auth_1,auth_2 = auth.split(' ');
        changeAuth(slide,auth_1,auth_2);
        presentation.save(f'./tmp/{name}_certificate.pptx');

        jpgconverter.convertToJpgFromCloud(f'./tmp/{name}_certificate.pptx');
        url = upload_certificate(f'./tmp/{name}_certificate.jpg');
        sendingmail.gmail_send_mail_with_attachment(
            From="Destroyer",
            TO=mail,
            name=name,
            file_name=f'./tmp/{name}_certificate.jpg',
            url=url
        );

        os.remove(f'./tmp/{name}_certificate.pptx');
        os.remove(f'./tmp/{name}_certificate.jpg');